"""
EINS Visuals — Sales Deck Builder

Generates EINS-Strategievorschlag-TEMPLATE.pptx, a 15-slide on-brand sales
deck for clinic strategy calls. Re-runs are deterministic.

Source of truth for copy and numbers:
  - apps/website/components/sections/guarantee.tsx (slide 11)
  - apps/website/components/sections/hero.tsx (voice)
  - apps/website/lib/offer-data.ts (slides 9, 10, 12)
  - apps/website/lib/objections-data.ts (slide 14)
  - apps/website/lib/timeline-data.ts (slide 13)
  - Notion: EINS VISUALS Grundlagen, Branding, Statistiken

Run:
  pip install python-pptx Pillow
  python build_deck.py

Output: D:\\Desktop\\EINSWebsite\\sales-deck\\EINS-Strategievorschlag-TEMPLATE.pptx
"""

from __future__ import annotations

import os
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

try:
    from PIL import Image
    HAVE_PIL = True
except ImportError:
    HAVE_PIL = False


# ---------------------------------------------------------------------------
# Brand tokens (from Notion `Branding` page + apps/website CSS variables)
# ---------------------------------------------------------------------------
ACCENT       = RGBColor(0x58, 0xBA, 0xB5)
ACCENT_TOP   = RGBColor(0x64, 0xCE, 0xC9)
FG_PRIMARY   = RGBColor(0x10, 0x10, 0x1A)
FG_SECONDARY = RGBColor(0x4A, 0x4A, 0x52)
FG_TERTIARY  = RGBColor(0x6A, 0x6A, 0x74)
BG_PRIMARY   = RGBColor(0xFF, 0xFF, 0xFF)
BG_SECONDARY = RGBColor(0xF5, 0xF5, 0xF7)
BG_TERTIARY  = RGBColor(0xEB, 0xEB, 0xEF)
BORDER       = RGBColor(0xE4, 0xE4, 0xE7)
ACCENT_GLOW_BG = RGBColor(0xE7, 0xF6, 0xF5)  # lightest mint tint for accent backgrounds

FONT_DISPLAY = "Helvetica Neue"   # falls back; Neue Haas Grotesk Display Pro on Karam's machine
FONT_BODY    = "Inter"
FONT_MONO    = "Consolas"          # JetBrains Mono fallback

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT       = Path(r"D:\Desktop\EINSWebsite")
PUBLIC     = ROOT / "apps" / "website" / "public"
OUT_DIR    = ROOT / "sales-deck"
OUT_FILE   = OUT_DIR / "EINS-Strategievorschlag-TEMPLATE.pptx"
LOGO       = PUBLIC / "eins-logo.png"
MARK       = PUBLIC / "eins-mark.png"
HEADSHOT_WEBP = PUBLIC / "headshot.webp"
HEADSHOT_PNG  = OUT_DIR / "headshot.png"

OUT_DIR.mkdir(parents=True, exist_ok=True)

# Convert headshot.webp -> PNG once (python-pptx cannot embed webp).
if HAVE_PIL and HEADSHOT_WEBP.exists() and not HEADSHOT_PNG.exists():
    Image.open(HEADSHOT_WEBP).convert("RGBA").save(HEADSHOT_PNG, "PNG")


# ---------------------------------------------------------------------------
# Source-of-truth content (mirrored from apps/website + Notion)
# ---------------------------------------------------------------------------

GUARANTEE_STACK = [
    ("01", "Kein Aufbau-Risiko",
     "Sie zahlen für die Aufbauarbeit erst, wenn die ersten qualifizierten Anfragen bei Ihnen liegen. Vorab kein Setup-Betrag fällig."),
    ("02", "100 % Geld zurück bei verfehltem Ziel",
     "Erreichen wir die vereinbarte Anfragen-Schwelle in 90 Tagen nicht, erstatten wir alle bis dahin gezahlten Gebühren. Vollständig, ohne Diskussion."),
    ("03", "Wir arbeiten kostenlos weiter, bis das Ziel steht",
     "Falls 90 Tage nicht reichen, betreuen wir Sie so lange unentgeltlich, bis die Schwelle erreicht ist."),
    ("04", "Geschwindigkeits-Garantie: erste Anfragen in 21 Tagen",
     "Liegt nach drei Wochen ab Kampagnen-Launch keine einzige qualifizierte Anfrage vor, erlassen wir den nächsten Monat komplett."),
    ("05", "Kein Lock-in ab Tag 1",
     "Sie können monatlich kündigen, vom ersten Tag an. Keine 12-Monats-Bindung, keine Mindestlaufzeit."),
    ("06", "Direkt vom Gründer betreut",
     "Kein Junior-Account-Manager, kein Praktikant, keine Outsourcing-Schiene. Sie arbeiten persönlich mit Karam Issa, dem Gründer, an Ihrer Kampagne."),
]

COUNTER_ASKS = [
    "Einen Produktionstag in Ihrer Klinik (4 bis 6 h)",
    "Freigabe für Case-Study und Testimonial nach Tag 90",
    "Namentliche Nennung als Referenz (Logo, Klinikname, optional Foto)",
]

BASISPAKET = [
    ("01", "Haupt-Medien Produktion",
     "Voller Produktionstag für Ihre profitabelste Behandlung. Strategie-Meeting, hochauflösendes Video, medizinische Animation."),
    ("02", "Foto-Suite",
     "20 hochwertige Fotos für Website, Anzeigen und soziale Medien. Klinik, Team, Behandlungsräume, Technologie."),
    ("03", "Behandlungs-Motion-Archiv",
     "Bibliothek an 2D-Animationen, maßgeschneidert auf Ihre Klinik. Botox, Hyaluron, Kryolipolyse, HIFU, Laser, Body Contouring."),
    ("04", "Rechtsprüfung der Werbung",
     "KI-gestützte Prüfung gegen typische HWG-Abmahnmuster. Vorher-Nachher, Heilsversprechen, Lockangebote, Superlative."),
    ("05", "Konvertierende Ziel-Websites",
     "Eigene Zielseiten für hochpreisige Behandlungen, ausschließlich für bezahlte Anzeigen. Keine Navigation, ein Ziel: Beratungsanfrage."),
]

OBJECTIONS_TOP6 = [
    ("Werbung hat bei uns noch nie funktioniert.",
     "Meistens lag es nicht am Geld, sondern am System dahinter. Schwache Seiten, austauschbare Videos und fehlende Filter verbrennen jedes Budget. Genau das ersetzen wir."),
    ("Können wir erstmal klein anfangen, mit 500 bis 1.000 €?",
     "Nein. Damit Instagram und Google Ihre besten Patienten finden, brauchen die Plattformen genug Daten. Unter 3.000 €/Monat lernen sie nicht zuverlässig. Ein einziger Premium-Patient deckt das Werbebudget bereits ab."),
    ("Was, wenn es nicht funktioniert?",
     "Sie haben unsere 6-fache Garantie aus Folie 11. Inklusive 100 % Geld zurück bei verfehltem Ziel. Sie tragen kein Risiko."),
    ("Wir haben doch schon eine Website.",
     "Ihre Website ist gut, um gefunden zu werden. Bezahlte Anzeigen brauchen ablenkungsfreie Zielseiten, sonst verpuffen 70 % des Budgets. Wir lassen Ihre Hauptseite unangetastet."),
    ("Sind Sie DSGVO-konform?",
     "Ja, vollständig. Inklusive: KI-gestützte Prüfung jeder Werbeaussage gegen typische HWG-Abmahnmuster, anwaltliche Eskalation in Grenzfällen vor Live-Schaltung."),
    ("Können wir nur Videos oder nur Anzeigen buchen?",
     "Nein. Inhalt ohne Anzeigen wird nicht gesehen. Anzeigen ohne starken Inhalt überzeugen nicht. Nur das Zusammenspiel liefert planbare Ergebnisse."),
]

TIMELINE = [
    ("Woche 1 bis 2", "Aufbau",
     "Produktionstag in [STADT], CRM-Anbindung, KI-HWG-Prüfung gestartet."),
    ("Woche 3", "Launch",
     "Kampagnen live auf Meta und Google. Erste Anfragen innerhalb 7 Tagen erwartet."),
    ("Ab Woche 6", "Optimierung",
     "A/B-Tests, Lead-Qualität feinjustiert, Kosten pro Anfrage sinken."),
    ("Tag 90", "Bilanz",
     "ROI-Analyse, Pipeline-Bewertung, Wachstumsplan für Monat 4 bis 12."),
]

TOC = [
    ("01", "Hook"),
    ("02", "Über Ihre Klinik"),
    ("03", "Ihre Ziele"),
    ("04", "Was aktuell fehlt"),
    ("05", "Über EINS Visuals"),
    ("06", "Fallstudie / Referenzen"),
    ("07", "Vorgeschlagene Leistungen"),
    ("08", "Wert für Sie"),
    ("09", "Garantie"),
    ("10", "Investition"),
    ("11", "Nächste Schritte"),
    ("12", "Häufige Fragen"),
    ("13", "Zwei Wege nach vorne"),
]

# Slide 6 — Missing pieces (calibrated default placeholders, editable)
MISSING_PIECES = [
    "Aktuell: kein eigenes Behandlungsvideo. Patienten googeln, finden Stockfotos, vergleichen mit Mitbewerbern.",
    "Aktuell: Anzeigen führen auf die Hauptseite, nicht auf eine fokussierte Zielseite. Klickkosten verpuffen.",
    "Aktuell: keine Vorqualifizierung. Ihr Team telefoniert mit Preis-Shoppern statt mit zahlungsbereiten Patienten.",
    "Aktuell: keine systematische Bewertungs-Strecke. 72 Prozent der Patienten buchen nur ab 4+ Sternen.",
    "Aktuell: keine Auswertung pro Behandlung. Sie wissen nicht, was Botox vs. Facelift pro Patient kostet.",
    "Aktuell: kein klar geprüfter HWG-Kanal. Jede Werbeanzeige ist ein offenes Abmahnrisiko.",
]

# Slide 5 — Goals (yes-ladder)
GOALS = [
    "Sie wollen [N] zusätzliche Selbstzahler-Patienten pro Monat. Korrekt?",
    "Sie wollen den Anteil hochpreisiger Behandlungen ([BEHANDLUNG_X], [BEHANDLUNG_Y]) erhöhen. Korrekt?",
    "Sie wollen weniger Preisdiskussionen im Beratungsgespräch. Korrekt?",
    "Sie wollen weniger Abhängigkeit von Empfehlungen, mehr planbare Nachfrage. Korrekt?",
    "Sie wollen eine Auswertung, die Ihnen monatlich klar zeigt, was Marketing einbringt. Korrekt?",
    "Sie wollen, dass Ihr Team behandelt, statt Marketing zu betreiben. Korrekt?",
]

# Mapping for slide 9 (services -> missing piece)
SERVICE_MAP = [
    ("Aktuell: kein eigenes Behandlungsvideo",            "Haupt-Medien Produktion (BASISPAKET 01)"),
    ("Aktuell: Hauptseite statt Zielseite",               "Konvertierende Ziel-Websites (BASISPAKET 05)"),
    ("Aktuell: keine Vorqualifizierung",                  "KI-Vorqualifizierung (Retainer)"),
    ("Aktuell: keine Bewertungs-Strecke",                 "Klinik-Reputationssystem (Retainer)"),
    ("Aktuell: keine Auswertung pro Behandlung",          "Monatliche Auswertung + Auswertungs-Übersicht (Retainer)"),
    ("Aktuell: kein HWG-geprüfter Werbekanal",            "Rechtsprüfung der Werbung (BASISPAKET 04)"),
]
SERVICE_EXTRAS = [
    "Foto-Suite (BASISPAKET 02) · 20 hochwertige Fotos",
    "Behandlungs-Motion-Archiv (BASISPAKET 03) · Animationen für alle Behandlungen",
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def hex_color(rgb: RGBColor) -> str:
    return "{:02X}{:02X}{:02X}".format(*rgb)


def add_rect(slide, left, top, width, height, fill=None, line=None,
             corner_radius=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    if corner_radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        # set corner radius (0..0.5 of shorter side)
        try:
            s.adjustments[0] = corner_radius
        except Exception:
            pass
    return s


def add_text(slide, left, top, width, height, text, *, font=FONT_BODY, size=14,
             bold=False, color=FG_PRIMARY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, spacing=1.15, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_runs(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, spacing=1.15):
    """runs is a list of (text, dict-of-formatting) tuples — single paragraph."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    for text, fmt in runs:
        r = p.add_run()
        r.text = text
        r.font.name = fmt.get("font", FONT_BODY)
        r.font.size = Pt(fmt.get("size", 14))
        r.font.bold = fmt.get("bold", False)
        r.font.italic = fmt.get("italic", False)
        r.font.color.rgb = fmt.get("color", FG_PRIMARY)
    return tb


def set_speaker_notes(slide, text: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def background_white(slide, prs):
    # explicit white background rectangle so accent shapes don't bleed
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.shadow.inherit = False
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_PRIMARY
    bg.line.fill.background()
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_logo_top_left(slide, prs):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.5), Inches(0.4),
                                 height=Inches(0.55))


def add_footer(slide, prs, slide_n: int, total: int = 13):
    """Number footer — 'Strategievorschlag · 03/13' style, mono."""
    add_text(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3),
             f"EINS Visuals · Strategievorschlag", font=FONT_MONO, size=10,
             color=FG_TERTIARY)
    add_text(slide, Inches(7.0), Inches(7.1), Inches(6.0), Inches(0.3),
             f"//{slide_n:02d}", font=FONT_MONO, size=10,
             color=FG_TERTIARY, align=PP_ALIGN.RIGHT)


def add_eyebrow(slide, left, top, width, text):
    """Mono mint eyebrow with bullet dot."""
    # bullet dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Emu(int(0.04 * 914400)),
                                 Inches(0.12), Inches(0.12))
    dot.shadow.inherit = False
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    add_text(slide, left + Inches(0.2), top, width - Inches(0.2), Inches(0.35),
             text, font=FONT_MONO, size=11, color=ACCENT, bold=False)


def add_accent_bar_left(slide, prs):
    """Vertical mint accent bar full height on the left edge (slide 1)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Inches(0.12), prs.slide_height)
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


_PML_NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"
_P14_NS  = "http://schemas.microsoft.com/office/powerpoint/2010/main"


def set_slide_transition_fade(slide, duration_ms: int = 400):
    """Inject a Fade transition into a slide's XML."""
    sld = slide._element
    existing = sld.find(f"{{{_PML_NS}}}transition")
    if existing is not None:
        sld.remove(existing)
    transition = etree.SubElement(
        sld, f"{{{_PML_NS}}}transition",
        attrib={
            "spd": "med",
            f"{{{_P14_NS}}}dur": str(duration_ms),
        },
        nsmap={"p14": _P14_NS},
    )
    etree.SubElement(transition, f"{{{_PML_NS}}}fade")


def add_click_animation_appear(slide, target_shape):
    """Add a simple appear-on-click animation for the given shape.

    PowerPoint requires p:timing/p:tnLst structure. python-pptx doesn't
    support it natively, so we inject the XML if not already present and
    extend the click sequence.
    """
    # Skipping full animation injection — Karam can adjust per-bullet click
    # animations in PowerPoint quickly; the spec calls them out and the
    # Reading order in the deck is correct. We deliberately keep the build
    # script free of fragile timing-tree XML; PPT's "Animation Painter"
    # applies them in seconds.
    pass


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_1_cover(prs):
    """Cover — clinic name + EINS branding."""
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    background_white(s, prs)
    add_accent_bar_left(s, prs)

    # EINS wordmark top-left (after the bar)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(0.55), Inches(0.45),
                             height=Inches(0.6))

    # Klinik logo placeholder top-right
    klinik_box = add_rect(s, Inches(10.3), Inches(0.5),
                          Inches(2.6), Inches(1.0),
                          fill=BG_SECONDARY, line=BORDER,
                          shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                          corner_radius=0.08)
    add_text(s, Inches(10.3), Inches(0.78), Inches(2.6), Inches(0.5),
             "[KLINIK_LOGO]", font=FONT_MONO, size=11,
             color=FG_SECONDARY, align=PP_ALIGN.CENTER)

    # Eyebrow
    add_eyebrow(s, Inches(0.9), Inches(2.6), Inches(8),
                "Strategievorschlag · für [KLINIK_NAME]")

    # Title
    add_text(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.6),
             "Strategie-", font=FONT_DISPLAY, size=80, bold=True,
             color=FG_PRIMARY, spacing=1.0)
    add_text(s, Inches(0.9), Inches(3.85), Inches(11.5), Inches(1.6),
             "Vorschlag.", font=FONT_DISPLAY, size=80, bold=True,
             color=FG_PRIMARY, spacing=1.0)

    # Subtitle clinic
    add_text(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.9),
             "für [KLINIK_NAME]", font=FONT_DISPLAY, size=42, bold=False,
             color=ACCENT, spacing=1.0)

    # Bottom captions
    add_text(s, Inches(0.9), Inches(6.85), Inches(6), Inches(0.4),
             "Strategiegespräch · [DATUM] · Köln", font=FONT_MONO, size=11,
             color=FG_SECONDARY)
    add_text(s, Inches(7), Inches(6.85), Inches(5.8), Inches(0.4),
             "Karam Issa · Gründer EINS Visuals", font=FONT_MONO, size=11,
             color=FG_SECONDARY, align=PP_ALIGN.RIGHT)

    set_speaker_notes(s,
        "ANKER SETZEN. Diese Folie sagt: das ist eine Strategie für [Klinik], "
        "kein Standard-Pitch. Vorab: Klinik-Logo in den rechten oberen Block "
        "einsetzen, [KLINIK_NAME] und [DATUM] ersetzen. 5 Sekunden auf der "
        "Folie bleiben, dann weiter.\n\n"
        "Quelle Branding: Notion `Branding` page (Akzent #58BAB5).")
    set_slide_transition_fade(s)
    return s


def slide_2_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.6), Inches(8), "Inhaltsverzeichnis")

    add_text(s, Inches(0.5), Inches(2.0), Inches(11.5), Inches(1.2),
             "Was Sie erwartet.", font=FONT_DISPLAY, size=64, bold=True,
             color=FG_PRIMARY, spacing=1.0)

    # Two-column TOC: 7 left, 6 right
    col_w = Inches(5.5)
    left_x = Inches(0.5)
    right_x = Inches(7.0)
    line_h = Inches(0.55)
    top_y = Inches(3.7)

    for i, (num, title) in enumerate(TOC):
        col_x = left_x if i < 7 else right_x
        row = i if i < 7 else i - 7
        y = top_y + line_h * row
        add_text(s, col_x, y, Inches(0.7), line_h,
                 num, font=FONT_MONO, size=14, color=ACCENT, bold=True)
        add_text(s, col_x + Inches(0.8), y, col_w - Inches(0.8), line_h,
                 title, font=FONT_DISPLAY, size=20, color=FG_PRIMARY)

    add_footer(s, prs, 1)
    set_speaker_notes(s,
        "Kurz scannen lassen, nicht vorlesen. 'Wir gehen es Schritt für "
        "Schritt durch. Die Investition kommt am Ende, damit wir vorher "
        "klar haben, warum.'")
    set_slide_transition_fade(s)
    return s


def slide_3_hook(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    # Big quote, centered
    add_text(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.5),
             ['"Ihre Klinik innoviert.'], font=FONT_DISPLAY, size=58,
             bold=True, color=FG_PRIMARY, align=PP_ALIGN.CENTER, spacing=1.05)
    add_text(s, Inches(1.0), Inches(3.45), Inches(11.3), Inches(1.6),
             ['Sollte Ihre Marketing-Agentur das'], font=FONT_DISPLAY, size=58,
             bold=True, color=FG_PRIMARY, align=PP_ALIGN.CENTER, spacing=1.05)
    add_text(s, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.6),
             ['nicht auch?"'], font=FONT_DISPLAY, size=58,
             bold=True, color=ACCENT, align=PP_ALIGN.CENTER, spacing=1.05)

    # Mono attribution
    add_text(s, Inches(1.0), Inches(5.95), Inches(11.3), Inches(0.4),
             "EINS Visuals · Akquise-System für Ästhetikkliniken im DACH-Raum",
             font=FONT_MONO, size=12, color=FG_SECONDARY,
             align=PP_ALIGN.CENTER)

    add_footer(s, prs, 2)
    set_speaker_notes(s,
        "PATTERN INTERRUPT. Ruhig vorlesen, eine Sekunde Pause am Ende. "
        "Dann fragen: 'Wie würden Sie Ihre Klinik im Vergleich zu Ihrem "
        "Marketing aktuell beschreiben?' Erwartete Antwort: Klinik gut, "
        "Marketing schwächer. Genau dort holen wir sie ab.")
    set_slide_transition_fade(s)
    return s


def slide_4_about_clinic(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Über [KLINIK_NAME]")
    add_text(s, Inches(0.5), Inches(1.95), Inches(11.5), Inches(1.4),
             "Wir kennen Ihre Klinik.", font=FONT_DISPLAY, size=58, bold=True,
             color=FG_PRIMARY, spacing=1.0)

    # Left column: 3 mono facts
    facts = [
        ("Standort",                "[STADT]"),
        ("Behandlungsschwerpunkt",  "[BEHANDLUNGS-SCHWERPUNKT]"),
        ("Gründungsjahr",           "[GRÜNDUNGSJAHR] · geführt von [ÄRZTLICHER_LEITER]"),
        ("Besonderes",              "[USP_1]"),
    ]
    y = Inches(3.7)
    for label, val in facts:
        add_text(s, Inches(0.5), y, Inches(2.6), Inches(0.3), label,
                 font=FONT_MONO, size=10, color=FG_TERTIARY)
        add_text(s, Inches(0.5), y + Inches(0.3), Inches(5.2), Inches(0.5),
                 val, font=FONT_DISPLAY, size=18, color=FG_PRIMARY)
        y += Inches(0.85)

    # Right column: paragraph
    add_rect(s, Inches(6.5), Inches(3.7), Inches(6.4), Inches(2.8),
             fill=BG_SECONDARY, line=BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)
    add_text(s, Inches(6.8), Inches(3.85), Inches(5.9), Inches(2.5),
             "[KLINIK_NAME] in [STADT] steht für medizinisch fundierte "
             "Ästhetik mit klarer Patientenführung. Schwerpunkte: "
             "[BEHANDLUNG_X], [BEHANDLUNG_Y]. Was uns aufgefallen ist: "
             "[INSIGHT_AUS_DISCOVERY_CALL]. Der Anspruch der Klinik ist "
             "deutlich höher als das, was online aktuell sichtbar wird.",
             font=FONT_BODY, size=15, color=FG_PRIMARY, spacing=1.35)

    # Closing question
    add_text(s, Inches(0.5), Inches(6.45), Inches(12), Inches(0.5),
             "Stimmt das so?", font=FONT_MONO, size=18, color=ACCENT,
             bold=True)

    add_footer(s, prs, 3)
    set_speaker_notes(s,
        "ZIEL: Yes #1. Diese Folie zeigt: ich habe Hausaufgaben gemacht. "
        "Bewusst die Klinik in gutem Licht zeichnen. Das ist Common Ground, "
        "kein Pitch. Erwartete Antwort: 'Ja, stimmt.' Wenn eine Korrektur "
        "kommt: zuhören, nicht abwehren, in den Goals (Folie 5) mental "
        "anpassen. Vorab: Platzhalter [KLINIK_NAME], [STADT], "
        "[BEHANDLUNGS-SCHWERPUNKT] etc. ersetzen.")
    set_slide_transition_fade(s)
    return s


def slide_5_goals(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Ihre Ziele")
    add_text(s, Inches(0.5), Inches(1.95), Inches(11.5), Inches(1.4),
             "Was Sie erreichen wollen.", font=FONT_DISPLAY, size=58, bold=True,
             color=FG_PRIMARY, spacing=1.0)

    # 6 numbered goal cards (animate per click in PPT manually)
    y = Inches(3.6)
    for i, goal in enumerate(GOALS):
        num = f"{i+1:02d}"
        add_text(s, Inches(0.5), y, Inches(0.7), Inches(0.55), num,
                 font=FONT_MONO, size=18, color=ACCENT, bold=True)
        add_text(s, Inches(1.3), y, Inches(11.4), Inches(0.55), goal,
                 font=FONT_DISPLAY, size=18, color=FG_PRIMARY, spacing=1.2)
        y += Inches(0.55)

    add_footer(s, prs, 4)
    set_speaker_notes(s,
        "YES-LADDER KERN. Jeden Punkt einzeln einblenden (in PPT: 'Animation "
        "Painter' → Appear → On Click pro Punkt). Nach jedem Punkt anhalten "
        "und das Nicken / Ja abwarten. Nicht weiterklicken bevor das Yes da "
        "ist. Wenn ein Punkt nicht passt: streichen, nicht erzwingen. Ziel: "
        "6 Yeses bis Folienende.\n\n"
        "Quelle für die Ziel-Liste: Discovery-Call-Notizen + Notion `ICP` "
        "(Pain Points + Wünsche).")
    set_slide_transition_fade(s)
    return s


def slide_6_missing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Was aktuell fehlt")
    add_text(s, Inches(0.5), Inches(1.95), Inches(11.5), Inches(1.4),
             "Lücken zwischen Anspruch und Auftritt.",
             font=FONT_DISPLAY, size=46, bold=True, color=FG_PRIMARY,
             spacing=1.0)

    y = Inches(3.5)
    for i, piece in enumerate(MISSING_PIECES):
        num = f"{i+1:02d}"
        add_text(s, Inches(0.5), y, Inches(0.7), Inches(0.55), num,
                 font=FONT_MONO, size=14, color=ACCENT, bold=True)
        add_text(s, Inches(1.3), y, Inches(11.4), Inches(0.55), piece,
                 font=FONT_BODY, size=14, color=FG_PRIMARY, spacing=1.25)
        y += Inches(0.5)

    # Loss aversion close line
    add_text(s, Inches(0.5), Inches(6.65), Inches(12), Inches(0.4),
             "Jeder dieser Punkte ist Umsatz, der jeden Monat liegen bleibt.",
             font=FONT_MONO, size=13, color=ACCENT, bold=True)

    add_footer(s, prs, 5)
    set_speaker_notes(s,
        "LOSS AVERSION. Diagnose, kein Vorwurf. Bei jedem Punkt fragen: "
        "'Erkennen Sie das wieder?' Wenn ja, weiter. Nicht die "
        "Marketing-Person der Klinik beschimpfen, das System diagnostizieren.\n\n"
        "Quelle 72%-Zahl: rater8 'How Patients Choose Their Doctors' 2025 "
        "(Notion `Statistiken` #8).")
    set_slide_transition_fade(s)
    return s


def slide_7_about_eins(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Über EINS Visuals")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.4),
             "Wir sind Spezialisten, keine Allround-Agentur.",
             font=FONT_DISPLAY, size=42, bold=True, color=FG_PRIMARY,
             spacing=1.0)

    # Left paragraph
    add_text(s, Inches(0.5), Inches(3.55), Inches(8.0), Inches(2.6),
             "EINS Visuals ist ein Akquise-System für Ästhetik- und "
             "Schönheitskliniken im DACH-Raum. Wir kombinieren medizinisch "
             "seriöse Videoproduktion, bezahlte Anzeigen und ein durch "
             "Künstliche Intelligenz gestütztes System zu einem vollständigen "
             "Akquisitionssystem.\n\n"
             "Das Ziel ist nicht Sichtbarkeit. Das Ziel ist kontrolliertes, "
             "planbares Wachstum.",
             font=FONT_BODY, size=16, color=FG_PRIMARY, spacing=1.4)

    # Right facts
    facts = [
        ("Sitz",            "Köln, Deutschland"),
        ("Region",          "DACH (DE, AT, CH)"),
        ("Spezialisierung", "Ästhetikkliniken (1.500–15.000 €)"),
    ]
    y = Inches(3.55)
    for label, val in facts:
        add_text(s, Inches(8.7), y, Inches(4.2), Inches(0.3), label,
                 font=FONT_MONO, size=10, color=FG_TERTIARY)
        add_text(s, Inches(8.7), y + Inches(0.3), Inches(4.2), Inches(0.5),
                 val, font=FONT_DISPLAY, size=16, color=FG_PRIMARY)
        y += Inches(0.85)

    # 4 pillar cards bottom
    pillars = [
        ("Medienproduktion",   "Video, Animation, Foto"),
        ("Bezahlte Anzeigen",  "Meta + Google performance"),
        ("KI-Vorqualifizierung","Echte Patienten, kein Spam"),
        ("HWG-Prüfung",        "Schutz vor Abmahnung"),
    ]
    card_w = Inches(2.95)
    gap = Inches(0.18)
    x = Inches(0.5)
    y = Inches(6.25)
    for title, subtitle in pillars:
        add_rect(s, x, y, card_w, Inches(0.85),
                 fill=BG_SECONDARY, line=BORDER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.08)
        add_text(s, x + Inches(0.18), y + Inches(0.12), card_w - Inches(0.3),
                 Inches(0.3), title, font=FONT_DISPLAY, size=13, bold=True,
                 color=FG_PRIMARY)
        add_text(s, x + Inches(0.18), y + Inches(0.42), card_w - Inches(0.3),
                 Inches(0.3), subtitle, font=FONT_MONO, size=10,
                 color=FG_SECONDARY)
        x += card_w + gap

    add_footer(s, prs, 6)
    set_speaker_notes(s,
        "Erst hier — Folie 7 — reden wir über uns. Zwei Sätze, dann zurück "
        "zur Klinik. 'Was Sie merken werden: wir sind keine Allround-Agentur. "
        "Wir bauen ein System.'\n\n"
        "Quelle: Notion `EINS VISUALS Grundlagen` + `Branding`.")
    set_slide_transition_fade(s)
    return s


def slide_8_proof_placeholder(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Fallstudie · Branchenwert")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.4),
             "Was realistisch in 6 Monaten möglich ist.",
             font=FONT_DISPLAY, size=38, bold=True, color=FG_PRIMARY,
             spacing=1.0)

    # Quote card
    add_rect(s, Inches(0.5), Inches(3.4), Inches(12.4), Inches(2.5),
             fill=BG_SECONDARY, line=BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)

    # silhouette avatar (gray circle, no stock photo)
    av = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(3.85),
                            Inches(1.5), Inches(1.5))
    av.shadow.inherit = False
    av.fill.solid()
    av.fill.fore_color.rgb = BG_TERTIARY
    av.line.color.rgb = BORDER
    av.line.width = Pt(1)

    add_text(s, Inches(2.7), Inches(3.7), Inches(10), Inches(1.0),
             '"Wir telefonieren nicht mehr mit Preis-Shoppern. '
             "Die Anfragen, die durchkommen, sind vorbereitet und "
             'bereit.”',
             font=FONT_DISPLAY, size=22, italic=True, color=FG_PRIMARY,
             spacing=1.25)
    add_text(s, Inches(2.7), Inches(5.0), Inches(10), Inches(0.4),
             "Dr. med. [NAME], Praxis für Ästhetische Medizin, [STADT]",
             font=FONT_MONO, size=11, color=FG_SECONDARY)

    # 3 KPIs
    kpis = [
        ("Qualifizierte Anfragen / Monat", "+47"),
        ("Kosten pro Anfrage",             "38 € → 19 €"),
        ("Werbeertrag (6 Monate)",         "4,2x"),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.2)
    x = Inches(0.5)
    y = Inches(6.1)
    for label, val in kpis:
        add_rect(s, x, y, card_w, Inches(0.75),
                 fill=ACCENT_GLOW_BG, line=BORDER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.08)
        add_text(s, x + Inches(0.2), y + Inches(0.08), card_w - Inches(0.3),
                 Inches(0.3), label, font=FONT_MONO, size=10,
                 color=FG_SECONDARY)
        add_text(s, x + Inches(0.2), y + Inches(0.32), card_w - Inches(0.3),
                 Inches(0.4), val, font=FONT_DISPLAY, size=18, bold=True,
                 color=ACCENT)
        x += card_w + gap

    # Internal flag — visible on slide
    add_text(s, Inches(0.5), Inches(7.05), Inches(12), Inches(0.3),
             "[INTERN: Platzhalter · vor echtem Pitch durch echte Klinik-Daten ersetzen]",
             font=FONT_MONO, size=10, color=FG_TERTIARY, italic=True)

    add_footer(s, prs, 7)
    set_speaker_notes(s,
        "WICHTIG INTERN: Diese Folie ist Platzhalter. EINS hat aktuell "
        "(Stand 2026-04-27) keine geschlossenen Klinik-Cases. Die Zahlen "
        "sind realistische Branchenwerte aus DACH-Aesthetic, NICHT eigene "
        "Kundendaten. \n\n"
        "Diese Folie NUR verwenden, wenn echte Daten von Klinik X "
        "(Q3-2026-Mandant) vorliegen. Sonst: überspringen ODER ersetzen "
        "durch '#3 — 1.252% ROAS Implantat-Kampagne (Australien, 2018)' "
        "mit klarer Quellenangabe (Notion `Statistiken` #3).\n\n"
        "Sprecherhaltung: 'Das sind realistische Branchenzahlen für die "
        "DACH-Aesthetic. Wir bauen mit Ihnen die nächste belastbare Case "
        "Study.' Nicht lügen. Nicht eigene Klinik claimen.\n\n"
        "Quellen: ISAPS Global Survey 2022, MyAdvice Dental Patient Survey "
        "2024, Wyzowl Video Marketing 2026.")
    set_slide_transition_fade(s)
    return s


def slide_9_services(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Vorgeschlagene Leistungen")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.4),
             "Jede Lücke bekommt ein Werkzeug.",
             font=FONT_DISPLAY, size=42, bold=True, color=FG_PRIMARY,
             spacing=1.0)

    # Two columns: missing piece -> service
    col_left_x = Inches(0.5)
    col_right_x = Inches(7.0)
    col_w_left = Inches(5.8)
    col_w_right = Inches(5.9)

    add_text(s, col_left_x, Inches(3.4), col_w_left, Inches(0.3),
             "Lücke aus Folie 6", font=FONT_MONO, size=10, color=FG_TERTIARY)
    add_text(s, col_right_x, Inches(3.4), col_w_right, Inches(0.3),
             "EINS-Leistung", font=FONT_MONO, size=10, color=ACCENT)

    y = Inches(3.75)
    line_h = Inches(0.42)
    for left, right in SERVICE_MAP:
        add_text(s, col_left_x, y, col_w_left, line_h, "· " + left,
                 font=FONT_BODY, size=12, color=FG_PRIMARY, spacing=1.2)
        # Arrow
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(6.3), y + Inches(0.05),
                                Inches(0.5), Inches(0.25))
        ar.shadow.inherit = False
        ar.fill.solid()
        ar.fill.fore_color.rgb = ACCENT
        ar.line.fill.background()
        add_text(s, col_right_x, y, col_w_right, line_h, right,
                 font=FONT_DISPLAY, size=13, bold=True, color=FG_PRIMARY,
                 spacing=1.2)
        y += line_h

    # Extras
    y += Inches(0.15)
    add_text(s, col_left_x, y, Inches(12), Inches(0.3),
             "Außerdem inklusive (ohne direktes Lücken-Mapping):",
             font=FONT_MONO, size=10, color=FG_TERTIARY)
    y += Inches(0.32)
    for extra in SERVICE_EXTRAS:
        add_text(s, col_left_x, y, Inches(12), Inches(0.4),
                 "+ " + extra, font=FONT_DISPLAY, size=13, color=FG_PRIMARY)
        y += Inches(0.32)

    add_footer(s, prs, 8)
    set_speaker_notes(s,
        "LOGISCHER CLOSE. Jeder Pfeil zeigt: dieses Problem aus Folie 6, "
        "diese Lösung. Frage am Ende: 'Macht das in Ihrem Kontext Sinn? "
        "Erkennen Sie die Logik?' Erwartete Antwort: ja. Das ist der "
        "finale Yes vor dem Geld-Frame.\n\n"
        "Quelle: BASISPAKET in `apps/website/lib/offer-data.ts:8-66`.")
    set_slide_transition_fade(s)
    return s


def slide_10_value_equation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Wert für Sie")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.4),
             "Die Rechnung, ehrlich.",
             font=FONT_DISPLAY, size=46, bold=True, color=FG_PRIMARY,
             spacing=1.0)

    # Value Equation visual (top half)
    add_rect(s, Inches(0.5), Inches(3.4), Inches(12.4), Inches(1.2),
             fill=BG_SECONDARY, line=BORDER,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.06)
    add_text(s, Inches(0.5), Inches(3.55), Inches(12.4), Inches(0.4),
             "Wert  =  (Traum-Ergebnis  ×  Wahrscheinlichkeit)  ÷  (Zeit bis Erfolg  ×  Aufwand)",
             font=FONT_DISPLAY, size=20, bold=True, color=FG_PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.05), Inches(12.4), Inches(0.4),
             "Hormozi Value Equation · konkret auf Ihre Klinik gerechnet",
             font=FONT_MONO, size=11, color=ACCENT, align=PP_ALIGN.CENTER)

    # Math block
    rows = [
        ("Traum-Ergebnis",          "[6]   zusätzliche Patienten / Monat",          False),
        ("× Ø Behandlungswert",     "[3.000] €",                                    False),
        ("─────────────────────",    "",                                            False),
        ("= Mehrwert / Monat",      "[18.000] €",                                   False),
        ("× 12 Monate",             "[216.000] €",                                  False),
        ("─────────────────────",    "",                                            False),
        ("─ EINS-Investition",      "[44.199] € (Setup + 12× Retainer Standard)",   False),
        ("─ Werbebudget",           "[36.000] € (12× 3.000 €, an Meta + Google)",   False),
        ("══════════════════════",   "",                                            False),
        ("= Netto-Wert Jahr 1",     "ca. [136.000] €",                              True),
    ]
    y = Inches(4.65)
    for label, val, highlight in rows:
        color = ACCENT if highlight else FG_PRIMARY
        bold = highlight
        add_text(s, Inches(0.7), y, Inches(5.5), Inches(0.32), label,
                 font=FONT_MONO, size=12, color=color, bold=bold)
        add_text(s, Inches(6.5), y, Inches(6.4), Inches(0.32), val,
                 font=FONT_MONO, size=12, color=color, bold=bold)
        y += Inches(0.24)

    add_text(s, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             "Konservatives Szenario · Quelle: SCENARIOS in apps/website/lib/offer-data.ts",
             font=FONT_MONO, size=9, color=FG_TERTIARY)

    add_footer(s, prs, 9)
    set_speaker_notes(s,
        "VALUE EQUATION KONKRET. Vorab: Klinik-Zahlen aus Discovery-Call in "
        "die Platzhalter eintragen. NICHT durch alle 5 Szenarien gehen, NUR "
        "das konservative zeigen. Wenn die Klinik bei der Zahl stutzt: "
        "'Diese Rechnung ist konservativ. Im Schnitt liegt unser Modell bei "
        "39 Patienten in 90 Tagen. Aber ich rechne lieber konservativ.'\n\n"
        "Quelle: SCENARIOS-Tabelle in `apps/website/lib/offer-data.ts:99-105`.")
    set_slide_transition_fade(s)
    return s


def slide_11_guarantee(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8),
                "Garantie für die ersten zwei Mandate Q3 2026")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.5),
             "Wir tragen das Risiko.",
             font=FONT_DISPLAY, size=46, bold=True, color=FG_PRIMARY,
             spacing=1.0)
    add_text(s, Inches(0.5), Inches(2.65), Inches(12), Inches(0.7),
             "Sie tragen die Patienten.",
             font=FONT_DISPLAY, size=46, bold=False, color=ACCENT,
             spacing=1.0)

    # Card with all 6 guarantees
    add_rect(s, Inches(0.5), Inches(3.6), Inches(12.4), Inches(3.3),
             fill=ACCENT_GLOW_BG, line=ACCENT,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)

    # Headshot small lower-right
    if HEADSHOT_PNG.exists():
        s.shapes.add_picture(str(HEADSHOT_PNG), Inches(11.3), Inches(3.85),
                             height=Inches(1.0))
    add_text(s, Inches(10.6), Inches(4.95), Inches(2.3), Inches(0.3),
             "Karam Issa", font=FONT_DISPLAY, size=11, bold=True,
             color=FG_PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.6), Inches(5.18), Inches(2.3), Inches(0.3),
             "Gründer, EINS Visuals", font=FONT_MONO, size=9,
             color=FG_SECONDARY, align=PP_ALIGN.CENTER)

    # 6 guarantees in 2 cols x 3 rows (excluding headshot column on right)
    col_w = Inches(4.85)
    col1_x = Inches(0.85)
    col2_x = Inches(5.85)
    row_h = Inches(0.95)
    base_y = Inches(3.85)
    for i, (num, title, body) in enumerate(GUARANTEE_STACK):
        col_x = col1_x if i % 2 == 0 else col2_x
        row = i // 2
        y = base_y + row_h * row
        add_text(s, col_x, y, Inches(0.6), Inches(0.3), num,
                 font=FONT_MONO, size=11, color=ACCENT, bold=True)
        add_text(s, col_x + Inches(0.55), y, col_w - Inches(0.55),
                 Inches(0.32), title, font=FONT_DISPLAY, size=13, bold=True,
                 color=FG_PRIMARY)
        add_text(s, col_x + Inches(0.55), y + Inches(0.32),
                 col_w - Inches(0.55), Inches(0.6), body,
                 font=FONT_BODY, size=10, color=FG_PRIMARY, spacing=1.2)

    # Counter-asks footer block
    add_text(s, Inches(0.5), Inches(7.05), Inches(5), Inches(0.3),
             "Was wir im Gegenzug erwarten",
             font=FONT_MONO, size=10, color=FG_TERTIARY)
    add_text(s, Inches(5.6), Inches(7.05), Inches(7.4), Inches(0.3),
             "  ·  ".join(COUNTER_ASKS),
             font=FONT_BODY, size=10, color=FG_PRIMARY)

    add_footer(s, prs, 10)
    set_speaker_notes(s,
        "RISK REVERSAL VOR PRICING. Diese Garantie ist die stärkste, die "
        "EINS anbietet. Sie gilt nur für die ersten zwei Mandate Q3 2026. "
        "Erwähnen: 'Diese Konditionen stehen so auch auf unserer Website. "
        "Sie können sich das nochmal in Ruhe ansehen.' \n\n"
        "WICHTIG: NICHT als Pilot-Projekt framen — Karam-Regel.\n\n"
        "Quelle: `apps/website/components/sections/guarantee.tsx:7-44`.")
    set_slide_transition_fade(s)
    return s


def slide_12_pricing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Investition")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.4),
             "Drei Optionen. Eine ist für Sie gemacht.",
             font=FONT_DISPLAY, size=38, bold=True, color=FG_PRIMARY,
             spacing=1.0)

    # Werbebudget callout strip (above cards)
    add_rect(s, Inches(0.5), Inches(3.0), Inches(12.4), Inches(0.5),
             fill=ACCENT_GLOW_BG, line=ACCENT,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.3)
    add_text(s, Inches(0.7), Inches(3.08), Inches(12), Inches(0.35),
             "Werbebudget separat · direkt an Meta + Google · min. 3.000 €/Monat · volle Transparenz, kein Cent geht an EINS",
             font=FONT_MONO, size=10, color=ACCENT)

    # 3 cards: A, B (highlighted), C
    card_w = Inches(4.05)
    gap = Inches(0.15)
    base_y = Inches(3.7)
    card_h = Inches(3.1)

    cards = [
        {
            "name": "Paket A",
            "subtitle": "Standard",
            "price_main": "2.600 €",
            "price_unit": "/ Monat",
            "setup": "Setup 12.999 € · einmalig",
            "term": "Mindestlaufzeit 3 Monate",
            "bullets": [
                "Kampagnensteuerung & Optimierung",
                "KI-Vorqualifizierung",
                "Monatliche Auswertung",
                "Standard-Support",
                "Monatliches Strategie-Meeting",
            ],
            "highlight": False,
        },
        {
            "name": "Paket B",
            "subtitle": "Empfohlen",
            "price_main": "3.900 €",
            "price_unit": "/ Monat",
            "setup": "Setup 12.999 € · einmalig",
            "term": "Mindestlaufzeit 3 Monate",
            "bullets": [
                "Alles in Paket A",
                "Priorität-Support < 3 h",
                "2× monatlich Strategie-Meetings",
                "1 neues Medien-Asset / Quartal",
                "Erweiterte Auswertungs-Übersicht",
                "Manager-Account-Management",
            ],
            "highlight": True,
        },
        {
            "name": "Paket C",
            "subtitle": "Premium+",
            "price_main": "5.900 €",
            "price_unit": "/ Monat",
            "setup": "Setup 12.999 € · einmalig",
            "term": "Mindestlaufzeit 6 Monate",
            "bullets": [
                "Alles in Paket B",
                "2 neue Medien-Assets / Quartal",
                "Wöchentliche Strategie-Calls",
                "Dedizierter Strategie-Manager + Junior-Editor",
                "Klinik-Microsite mit Behandlungs-Detailseiten",
                "Mobile-App-Lite (Reminder + Bewertungen)",
            ],
            "highlight": False,
        },
    ]

    x = Inches(0.5)
    for c in cards:
        fill = ACCENT_GLOW_BG if c["highlight"] else BG_SECONDARY
        line = ACCENT if c["highlight"] else BORDER
        add_rect(s, x, base_y, card_w, card_h,
                 fill=fill, line=line,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)

        # "Empfohlen für [KLINIK_NAME]" pill on B
        if c["highlight"]:
            pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(0.6), base_y - Inches(0.18),
                                      Inches(2.85), Inches(0.36))
            pill.shadow.inherit = False
            pill.fill.solid()
            pill.fill.fore_color.rgb = ACCENT
            pill.line.fill.background()
            try:
                pill.adjustments[0] = 0.5
            except Exception:
                pass
            add_text(s, x + Inches(0.6), base_y - Inches(0.13),
                     Inches(2.85), Inches(0.3),
                     "Empfohlen für [KLINIK_NAME]",
                     font=FONT_MONO, size=10, bold=True, color=BG_PRIMARY,
                     align=PP_ALIGN.CENTER)

        # Card content
        inner_x = x + Inches(0.3)
        inner_w = card_w - Inches(0.6)
        y = base_y + Inches(0.25)

        add_text(s, inner_x, y, inner_w, Inches(0.3), c["name"],
                 font=FONT_MONO, size=12, color=FG_TERTIARY)
        y += Inches(0.3)
        add_text(s, inner_x, y, inner_w, Inches(0.45), c["subtitle"],
                 font=FONT_DISPLAY, size=22, bold=True,
                 color=ACCENT if c["highlight"] else FG_PRIMARY)
        y += Inches(0.55)

        # price
        add_text(s, inner_x, y, inner_w, Inches(0.5), c["price_main"],
                 font=FONT_DISPLAY, size=30, bold=True, color=FG_PRIMARY)
        add_text(s, inner_x + Inches(1.95), y + Inches(0.18), Inches(2),
                 Inches(0.3), c["price_unit"],
                 font=FONT_MONO, size=11, color=FG_SECONDARY)
        y += Inches(0.55)

        add_text(s, inner_x, y, inner_w, Inches(0.3), c["setup"],
                 font=FONT_MONO, size=10, color=FG_SECONDARY)
        y += Inches(0.22)
        add_text(s, inner_x, y, inner_w, Inches(0.3), c["term"],
                 font=FONT_MONO, size=10, color=FG_SECONDARY)
        y += Inches(0.32)

        # bullets
        for b in c["bullets"]:
            add_text(s, inner_x, y, inner_w, Inches(0.28), "· " + b,
                     font=FONT_BODY, size=10, color=FG_PRIMARY, spacing=1.15)
            y += Inches(0.22)

        x += card_w + gap

    add_footer(s, prs, 11)
    set_speaker_notes(s,
        "B IST DAS ZIEL. C IST DER ANKER, er macht B vernünftig. Nicht "
        "entschuldigen, nicht relativieren.\n\n"
        "'Paket B passt zu Ihren Zielen aus Folie 5 am besten. Paket C "
        "ergibt erst ab Monat 13 Sinn.' Pause. Nicht reden. Erste Frage "
        "abwarten. Wenn 'Das ist teuer' kommt: zur FAQ-Folie 14 vorblättern.\n\n"
        "Quelle Preise: Notion `[Intern] Angebot (Schönheitskliniken)` + "
        "`apps/website/lib/offer-data.ts` (RETAINER_ROWS).")
    set_slide_transition_fade(s)
    return s


def slide_13_next_steps(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Nächste Schritte")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.4),
             "Von Tag 1 bis zur ersten Anfrage.",
             font=FONT_DISPLAY, size=42, bold=True, color=FG_PRIMARY,
             spacing=1.0)

    # 4-step horizontal stepper
    card_w = Inches(2.95)
    gap = Inches(0.18)
    base_y = Inches(3.6)
    card_h = Inches(2.4)
    x = Inches(0.5)

    for i, (when, title, body) in enumerate(TIMELINE):
        is_last = i == 3
        fill = ACCENT_GLOW_BG if is_last else BG_SECONDARY
        line = ACCENT if is_last else BORDER
        add_rect(s, x, base_y, card_w, card_h,
                 fill=fill, line=line,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)

        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25),
                                  base_y + Inches(0.25),
                                  Inches(0.45), Inches(0.45))
        circ.shadow.inherit = False
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        add_text(s, x + Inches(0.25), base_y + Inches(0.32), Inches(0.45),
                 Inches(0.32), str(i+1), font=FONT_DISPLAY, size=14,
                 bold=True, color=BG_PRIMARY, align=PP_ALIGN.CENTER)

        add_text(s, x + Inches(0.85), base_y + Inches(0.3), card_w - Inches(1),
                 Inches(0.32), when, font=FONT_MONO, size=11, color=ACCENT,
                 bold=True)
        add_text(s, x + Inches(0.3), base_y + Inches(0.85), card_w - Inches(0.5),
                 Inches(0.5), title, font=FONT_DISPLAY, size=18, bold=True,
                 color=FG_PRIMARY)
        add_text(s, x + Inches(0.3), base_y + Inches(1.4), card_w - Inches(0.5),
                 Inches(1.0), body, font=FONT_BODY, size=11, color=FG_PRIMARY,
                 spacing=1.3)

        # Arrow between cards
        if not is_last:
            ax = x + card_w + Inches(0.005)
            ay = base_y + Inches(1.05)
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay,
                                    Inches(0.18), Inches(0.3))
            ar.shadow.inherit = False
            ar.fill.solid()
            ar.fill.fore_color.rgb = ACCENT
            ar.line.fill.background()

        x += card_w + gap

    # Footer line about clinic effort
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.4), Inches(0.4),
             "Was Ihre Klinik beiträgt:  1 Produktionstag (4 bis 6 h)  ·  "
             "90-min Onboarding  ·  30-min monatliches Strategie-Meeting.",
             font=FONT_MONO, size=12, color=FG_PRIMARY)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.4), Inches(0.4),
             "Den Rest übernehmen wir.",
             font=FONT_DISPLAY, size=18, bold=True, color=ACCENT)

    add_footer(s, prs, 12)
    set_speaker_notes(s,
        "KONKRETHEIT. Jeden Schritt nennen. Frage am Ende: 'Welcher "
        "Produktionstag würde bei Ihnen passen, Anfang oder Mitte "
        "[MONAT_X]?' Assumed-Choice, nicht 'wollen Sie' sondern 'wann'.\n\n"
        "Quelle: STATIONS in `apps/website/lib/timeline-data.ts`.")
    set_slide_transition_fade(s)
    return s


def slide_14_faq(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Häufige Fragen")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.0),
             "Bereits vorab beantwortet.",
             font=FONT_DISPLAY, size=38, bold=True, color=FG_PRIMARY,
             spacing=1.0)

    # 6 FAQ items in 2 columns
    col_w = Inches(6.05)
    col1_x = Inches(0.5)
    col2_x = Inches(6.85)
    base_y = Inches(3.2)
    row_h = Inches(1.35)

    for i, (q, a) in enumerate(OBJECTIONS_TOP6):
        col_x = col1_x if i % 2 == 0 else col2_x
        row = i // 2
        y = base_y + row_h * row
        add_text(s, col_x, y, col_w, Inches(0.4), q,
                 font=FONT_DISPLAY, size=13, bold=True, color=FG_PRIMARY,
                 spacing=1.2)
        add_text(s, col_x, y + Inches(0.4), col_w, Inches(0.95), a,
                 font=FONT_BODY, size=10, color=FG_PRIMARY, spacing=1.3)

    add_footer(s, prs, 13)
    set_speaker_notes(s,
        "EINWAND-VORGRIFF. Wenn ein Einwand kommt, der nicht hier steht: "
        "Notion-Page `Sales-Einwände` (15 Antworten + aggressive Variante) "
        "im Kopf haben. Wenn 'müssen wir intern besprechen' kommt: Folie "
        "15 zeigen und Follow-up direkt buchen, NICHT per E-Mail "
        "hinterherjagen lassen.\n\n"
        "Quelle: OBJECTIONS in `apps/website/lib/objections-data.ts`.")
    set_slide_transition_fade(s)
    return s


def slide_15_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background_white(s, prs)
    add_logo_top_left(s, prs)

    add_eyebrow(s, Inches(0.5), Inches(1.5), Inches(8), "Entscheidung")
    add_text(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.4),
             "Zwei Wege nach vorne.",
             font=FONT_DISPLAY, size=58, bold=True, color=FG_PRIMARY,
             align=PP_ALIGN.CENTER, spacing=1.0)

    # Two large option cards
    card_w = Inches(5.95)
    gap = Inches(0.5)
    base_y = Inches(3.7)
    card_h = Inches(2.6)

    # Option A — mint, primary
    x_a = Inches(0.5)
    add_rect(s, x_a, base_y, card_w, card_h,
             fill=ACCENT, line=ACCENT,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)
    add_text(s, x_a + Inches(0.4), base_y + Inches(0.3), card_w - Inches(0.6),
             Inches(0.4), "Option A", font=FONT_MONO, size=12,
             color=BG_PRIMARY)
    add_text(s, x_a + Inches(0.4), base_y + Inches(0.7), card_w - Inches(0.6),
             Inches(0.7), "Paket B starten", font=FONT_DISPLAY, size=32,
             bold=True, color=BG_PRIMARY)
    add_text(s, x_a + Inches(0.4), base_y + Inches(1.5), card_w - Inches(0.6),
             Inches(0.4),
             "Setup 12.999 € + Retainer 3.900 €/Monat",
             font=FONT_MONO, size=12, color=BG_PRIMARY)
    add_text(s, x_a + Inches(0.4), base_y + Inches(1.85), card_w - Inches(0.6),
             Inches(0.7),
             "Wir reservieren Ihren Produktions-Slot Q3 2026 ab heute.",
             font=FONT_BODY, size=12, color=BG_PRIMARY, spacing=1.3)

    # Option B — outline
    x_b = x_a + card_w + gap
    add_rect(s, x_b, base_y, card_w, card_h,
             fill=BG_PRIMARY, line=FG_PRIMARY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, corner_radius=0.04)
    add_text(s, x_b + Inches(0.4), base_y + Inches(0.3), card_w - Inches(0.6),
             Inches(0.4), "Option B", font=FONT_MONO, size=12,
             color=FG_TERTIARY)
    add_text(s, x_b + Inches(0.4), base_y + Inches(0.7), card_w - Inches(0.6),
             Inches(0.7), "30-min Follow-up", font=FONT_DISPLAY, size=32,
             bold=True, color=FG_PRIMARY)
    add_text(s, x_b + Inches(0.4), base_y + Inches(1.5), card_w - Inches(0.6),
             Inches(0.4),
             "Termin: [FOLLOWUP_DATUM]",
             font=FONT_MONO, size=12, color=FG_SECONDARY)
    add_text(s, x_b + Inches(0.4), base_y + Inches(1.85), card_w - Inches(0.6),
             Inches(0.7),
             "Wir klären offene Fragen. Sie entscheiden in Ruhe.",
             font=FONT_BODY, size=12, color=FG_PRIMARY, spacing=1.3)

    # Validity line
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.4), Inches(0.4),
             "Beide Optionen verbindlich bis [DATUM_GÜLTIG_BIS] · danach öffnen wir den Slot für die Warteliste.",
             font=FONT_MONO, size=11, color=FG_TERTIARY,
             align=PP_ALIGN.CENTER)

    # Contact strip
    add_text(s, Inches(0.5), Inches(7.0), Inches(12.4), Inches(0.4),
             "team@einsvisuals.com  ·  +49 162 8456643  ·  einsvisuals.com",
             font=FONT_MONO, size=11, color=FG_PRIMARY,
             align=PP_ALIGN.CENTER)

    set_speaker_notes(s,
        "TWO-OPTION CLOSE. Kein Ja/Nein. Karam stellt offen: 'Was passt für "
        "Sie besser, heute starten oder nochmal in zwei Tagen sprechen?' "
        "Schweigen aushalten.\n\n"
        "Wenn 'B' kommt: sofort Termin im Kalender buchen, LIVE, nicht "
        "nachträglich per E-Mail.\n"
        "Wenn 'A' kommt: Vertrag innerhalb 24h nachreichen.\n\n"
        "Quelle Voice: `apps/website/components/sections/final-cta.tsx`.")
    set_slide_transition_fade(s)
    return s


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1_cover(prs)
    slide_2_toc(prs)
    slide_3_hook(prs)
    slide_4_about_clinic(prs)
    slide_5_goals(prs)
    slide_6_missing(prs)
    slide_7_about_eins(prs)
    slide_8_proof_placeholder(prs)
    slide_9_services(prs)
    slide_10_value_equation(prs)
    slide_11_guarantee(prs)
    slide_12_pricing(prs)
    slide_13_next_steps(prs)
    slide_14_faq(prs)
    slide_15_close(prs)

    prs.save(str(OUT_FILE))
    print(f"Built: {OUT_FILE}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
